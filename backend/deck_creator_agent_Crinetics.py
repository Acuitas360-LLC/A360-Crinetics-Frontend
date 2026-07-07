"""
ppt_theme_aware.py
==================
Premium slide generator that inherits theme (colors, fonts) from an uploaded
.pptx file and renders on a clean WHITE background.

Usage
-----
    from ppt_theme_aware import create_ppt, extract_theme_from_pptx

    theme = extract_theme_from_pptx("uploaded_deck.pptx")
    prs   = create_ppt(slide_data, theme=theme)
    prs.save("output.pptx")

If no theme is provided the original dark-navy defaults are used.
"""


from __future__ import annotations
import os
from datetime import date
from typing import Any
import plotly.io as pio
pio.templates.default = "plotly"
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

import pandas as pd
from langchain_core.messages import BaseMessage, HumanMessage, AIMessage
from dotenv import load_dotenv
load_dotenv()
import json
from openai import OpenAI
from langchain_core.messages import BaseMessage, HumanMessage, AIMessage
import datetime
import warnings
import uuid
from langchain_core.messages import HumanMessage, AIMessage
from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
warnings.filterwarnings("ignore")


messages=[HumanMessage(content='How are sales trending?', additional_kwargs={}, response_metadata={}, id='8823ba60-6697-49c3-9f96-728e740c1032'), AIMessage(content='SQL Query Executed:\nWITH max_week_end_date_cte AS (\n    SELECT MAX(week_end_date) AS max_week_end_date\n    FROM data_867\n),\nr13w_window_cte AS (\n    SELECT\n        DATEADD(WEEK, -12, max_week_end_date) AS window_start_date,\n        max_week_end_date AS window_end_date\n    FROM max_week_end_date_cte\n),\nweekly_trend_cte AS (\n    SELECT\n        DATEADD(DAY, -6, d.week_end_date) AS week_start_date,\n        d.week_end_date,\n        COUNT(DISTINCT CASE WHEN d.is_business_day = 1 THEN d.date END) AS weekly_business_day_count,\n        SUM(d.relmora_total_mg) AS weekly_sales,\n        CAST(\n            SUM(d.relmora_total_mg) / NULLIF(COUNT(DISTINCT CASE WHEN d.is_business_day = 1 THEN d.date END), 0)\n            AS INTEGER\n        ) AS weekly_daily_average_sales\n    FROM data_867 d\n    JOIN r13w_window_cte w\n      ON d.week_end_date >= w.window_start_date\n     AND d.week_end_date <= w.window_end_date\n    GROUP BY d.week_end_date\n),\ntotal_r13w_cte AS (\n    SELECT\n        SUM(d.relmora_total_mg) AS total_mg_r13w\n    FROM data_867 d\n    JOIN r13w_window_cte w\n      ON d.week_end_date >= w.window_start_date\n     AND d.week_end_date <= w.window_end_date\n)\nSELECT\n    wt.week_start_date,\n    wt.week_end_date,\n    wt.weekly_business_day_count,\n    wt.weekly_sales,\n    wt.weekly_daily_average_sales,\n    tr.total_mg_r13w\nFROM weekly_trend_cte wt\nCROSS JOIN total_r13w_cte tr\nORDER BY wt.week_end_date ASC;\n\nResult Summary:\nThe analysis answers a straightforward business question: how Relmora’s national sales have trended over the most recent 13 weeks of available data, looking at both total weekly volume and the average sold per business day. Over this period, Relmora generated 1,054,069 mg in total national sales, with weekly results generally ranging from about 58,000 mg to 116,000 mg, showing meaningful week-to-week volatility rather than a steady trend. The strongest week was Jan 31–Feb 6 at 116,466 mg, while holiday-affected and shortened weeks were noticeably softer, including Dec 20–26 at 58,045 mg and Feb 14–20 at 67,351 mg.\n\nWhat stands out is that business-day-adjusted performance tells a more balanced story than raw weekly totals alone. Several 4-day weeks still posted relatively strong daily averages, suggesting demand held up better than the lower weekly totals might imply. The latest week, Feb 28–Mar 6, shows only 1 business day and 24,205 mg in sales, which indicates the most recent period is incomplete and should not be compared directly with prior full weeks. Overall, the 13-week view points to a pattern of fluctuating but resilient national demand, with timing effects from shortened weeks materially influencing the headline weekly totals.\n\nRelevant Questions:\n- What is the sales in the recent month?\n- How has our sales performance evolved over the past year?\n- Are we seeing strong short-term sales momentum?\n- How are we doing in terms of adding new businesses?', additional_kwargs={}, response_metadata={}, id='4c1271ba-c619-49e4-b3c0-dd31be76827a', tool_calls=[], invalid_tool_calls=[]), AIMessage(content='SQL query results', additional_kwargs={'type': 'sql_result', 'data': {'columns': ['WEEK_START_DATE', 'WEEK_END_DATE', 'WEEKLY_BUSINESS_DAY_COUNT', 'WEEKLY_SALES', 'WEEKLY_DAILY_AVERAGE_SALES', 'TOTAL_MG_R13W'], 'data': [{'WEEK_START_DATE': datetime.date(2025, 12, 6), 'WEEK_END_DATE': datetime.date(2025, 12, 12), 'WEEKLY_BUSINESS_DAY_COUNT': 5, 'WEEKLY_SALES': 84929, 'WEEKLY_DAILY_AVERAGE_SALES': 16986, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2025, 12, 13), 'WEEK_END_DATE': datetime.date(2025, 12, 19), 'WEEKLY_BUSINESS_DAY_COUNT': 5, 'WEEKLY_SALES': 95833, 'WEEKLY_DAILY_AVERAGE_SALES': 19167, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2025, 12, 20), 'WEEK_END_DATE': datetime.date(2025, 12, 26), 'WEEKLY_BUSINESS_DAY_COUNT': 4, 'WEEKLY_SALES': 58045, 'WEEKLY_DAILY_AVERAGE_SALES': 14511, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2025, 12, 27), 'WEEK_END_DATE': datetime.date(2026, 1, 2), 'WEEKLY_BUSINESS_DAY_COUNT': 4, 'WEEKLY_SALES': 77127, 'WEEKLY_DAILY_AVERAGE_SALES': 19282, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2026, 1, 3), 'WEEK_END_DATE': datetime.date(2026, 1, 9), 'WEEKLY_BUSINESS_DAY_COUNT': 5, 'WEEKLY_SALES': 108805, 'WEEKLY_DAILY_AVERAGE_SALES': 21761, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2026, 1, 10), 'WEEK_END_DATE': datetime.date(2026, 1, 16), 'WEEKLY_BUSINESS_DAY_COUNT': 5, 'WEEKLY_SALES': 90146, 'WEEKLY_DAILY_AVERAGE_SALES': 18029, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2026, 1, 17), 'WEEK_END_DATE': datetime.date(2026, 1, 23), 'WEEKLY_BUSINESS_DAY_COUNT': 4, 'WEEKLY_SALES': 84130, 'WEEKLY_DAILY_AVERAGE_SALES': 21033, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2026, 1, 24), 'WEEK_END_DATE': datetime.date(2026, 1, 30), 'WEEKLY_BUSINESS_DAY_COUNT': 5, 'WEEKLY_SALES': 79383, 'WEEKLY_DAILY_AVERAGE_SALES': 15877, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2026, 1, 31), 'WEEK_END_DATE': datetime.date(2026, 2, 6), 'WEEKLY_BUSINESS_DAY_COUNT': 5, 'WEEKLY_SALES': 116466, 'WEEKLY_DAILY_AVERAGE_SALES': 23293, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2026, 2, 7), 'WEEK_END_DATE': datetime.date(2026, 2, 13), 'WEEKLY_BUSINESS_DAY_COUNT': 5, 'WEEKLY_SALES': 81357, 'WEEKLY_DAILY_AVERAGE_SALES': 16271, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2026, 2, 14), 'WEEK_END_DATE': datetime.date(2026, 2, 20), 'WEEKLY_BUSINESS_DAY_COUNT': 4, 'WEEKLY_SALES': 67351, 'WEEKLY_DAILY_AVERAGE_SALES': 16838, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2026, 2, 21), 'WEEK_END_DATE': datetime.date(2026, 2, 27), 'WEEKLY_BUSINESS_DAY_COUNT': 5, 'WEEKLY_SALES': 86292, 'WEEKLY_DAILY_AVERAGE_SALES': 17258, 'TOTAL_MG_R13W': 1054069}, {'WEEK_START_DATE': datetime.date(2026, 2, 28), 'WEEK_END_DATE': datetime.date(2026, 3, 6), 'WEEKLY_BUSINESS_DAY_COUNT': 1, 'WEEKLY_SALES': 24205, 'WEEKLY_DAILY_AVERAGE_SALES': 24205, 'TOTAL_MG_R13W': 1054069}]}}, response_metadata={}, id='82e56320-9b1d-426c-a824-58e44e5072db', tool_calls=[], invalid_tool_calls=[]), AIMessage(content='Visualization', additional_kwargs={'type': 'visualization', 'code': 'import pandas as pd\nimport plotly.graph_objects as go\nfrom plotly.subplots import make_subplots\n\ndf = df.copy()\ndf["WEEK_END_DATE"] = pd.to_datetime(df["WEEK_END_DATE"])\ndf["WEEK_START_DATE"] = pd.to_datetime(df["WEEK_START_DATE"])\ndf = df.sort_values("WEEK_END_DATE")\n\nx_labels = df["WEEK_END_DATE"].dt.strftime("%Y-%m-%d")\n\nfig = make_subplots(specs=[[{"secondary_y": True}]])\n\nfig.add_trace(\n    go.Bar(\n        x=x_labels,\n        y=df["WEEKLY_SALES"],\n        name="Weekly Sales",\n        marker_color="#4C78A8",\n        opacity=0.85,\n        customdata=list(\n            zip(\n                df["WEEK_START_DATE"].dt.strftime("%Y-%m-%d"),\n                df["WEEK_END_DATE"].dt.strftime("%Y-%m-%d"),\n                df["WEEKLY_BUSINESS_DAY_COUNT"],\n                df["WEEKLY_DAILY_AVERAGE_SALES"]\n            )\n        ),\n        hovertemplate=f"Week Start: %{{customdata[0]}}<br>Week End: %{{customdata[1]}}<br>Weekly Sales: %{{y:,.0f}}<br>Business Days: %{{customdata[2]}}<br>Daily Avg Sales: %{{customdata[3]:,.0f}}<extra></extra>"\n    ),\n    secondary_y=False\n)\n\nfig.add_trace(\n    go.Scatter(\n        x=x_labels,\n        y=df["WEEKLY_DAILY_AVERAGE_SALES"],\n        name="Weekly Daily Avg Sales",\n        mode="lines+markers",\n        line=dict(color="#F58518", width=3),\n        marker=dict(size=8),\n        customdata=list(\n            zip(\n                df["WEEK_START_DATE"].dt.strftime("%Y-%m-%d"),\n                df["WEEK_END_DATE"].dt.strftime("%Y-%m-%d"),\n                df["WEEKLY_BUSINESS_DAY_COUNT"],\n                df["WEEKLY_SALES"]\n            )\n        ),\n        hovertemplate=f"Week Start: %{{customdata[0]}}<br>Week End: %{{customdata[1]}}<br>Daily Avg Sales: %{{y:,.0f}}<br>Business Days: %{{customdata[2]}}<br>Weekly Sales: %{{customdata[3]:,.0f}}<extra></extra>"\n    ),\n    secondary_y=True\n)\n\nlatest = df.iloc[-1]\nfig.add_trace(\n    go.Scatter(\n        x=[latest["WEEK_END_DATE"].strftime("%Y-%m-%d")],\n        y=[latest["WEEKLY_DAILY_AVERAGE_SALES"]],\n        name="Latest Week",\n        mode="markers",\n        marker=dict(color="#E45756", size=12, symbol="diamond"),\n        customdata=[[\n            latest["WEEK_START_DATE"].strftime("%Y-%m-%d"),\n            latest["WEEK_END_DATE"].strftime("%Y-%m-%d"),\n            latest["WEEKLY_BUSINESS_DAY_COUNT"],\n            latest["WEEKLY_SALES"]\n        ]],\n        hovertemplate=f"Latest Week<br>Week Start: %{{customdata[0]}}<br>Week End: %{{customdata[1]}}<br>Daily Avg Sales: %{{y:,.0f}}<br>Business Days: %{{customdata[2]}}<br>Weekly Sales: %{{customdata[3]:,.0f}}<extra></extra>",\n        showlegend=True\n    ),\n    secondary_y=True\n)\n\nfig.update_layout(\n    title=f"Relmora National Weekly Sales Trend - Most Recent 13 Weeks<br><sup>Total R13W Sales: {df[\'TOTAL_MG_R13W\'].iloc[0]:,.0f}</sup>",\n    xaxis_title="Week End Date",\n    yaxis_title="Weekly Sales",\n    yaxis2_title="Weekly Daily Average Sales",\n    template="plotly_white",\n    width=1000,\n    legend=dict(orientation="h", y=1.12, x=0),\n    hovermode="x unified"\n)\n\nfig.update_xaxes(type="category", tickangle=-45)\nfig.update_yaxes(tickformat=",", secondary_y=False)\nfig.update_yaxes(tickformat=",", secondary_y=True)'}, response_metadata={}, id='30c501ea-1993-4663-b307-b09a089d608a', tool_calls=[], invalid_tool_calls=[]), HumanMessage(content='Are we seeing strong short-term sales momentum?', additional_kwargs={}, response_metadata={}, id='bb36a2c0-349f-4e8a-b28b-1176c7ec4962'), AIMessage(content="SQL Query Executed:\nWITH max_week_end_date_cte AS (\n    SELECT MAX(week_end_date) AS max_week_end_date\n    FROM data_867\n),\nperiod_boundaries_cte AS (\n    SELECT\n        m.max_week_end_date AS r4w_end_date,\n        DATEADD(DAY, -27, m.max_week_end_date) AS r4w_start_date,\n        DATEADD(DAY, -28, m.max_week_end_date) AS p4w_end_date,\n        DATEADD(DAY, -55, m.max_week_end_date) AS p4w_start_date\n    FROM max_week_end_date_cte m\n),\nbusiness_day_counts_cte AS (\n    SELECT\n        COUNT(DISTINCT CASE WHEN d.is_business_day = 1 AND d.week_end_date BETWEEN pb.r4w_start_date AND pb.r4w_end_date THEN d.date END) AS r4w_business_days,\n        COUNT(DISTINCT CASE WHEN d.is_business_day = 1 AND d.week_end_date BETWEEN pb.p4w_start_date AND pb.p4w_end_date THEN d.date END) AS p4w_business_days\n    FROM data_867 d\n    CROSS JOIN period_boundaries_cte pb\n    WHERE d.week_end_date <= pb.r4w_end_date\n      AND d.week_end_date >= pb.p4w_start_date\n),\nperiod_completeness_cte AS (\n    SELECT\n        CASE WHEN COUNT(DISTINCT CASE WHEN d.week_end_date BETWEEN pb.r4w_start_date AND pb.r4w_end_date THEN d.week_end_date END) = 4 THEN 1 ELSE 0 END AS is_r4w_complete,\n        CASE WHEN COUNT(DISTINCT CASE WHEN d.week_end_date BETWEEN pb.p4w_start_date AND pb.p4w_end_date THEN d.week_end_date END) = 4 THEN 1 ELSE 0 END AS is_p4w_complete\n    FROM data_867 d\n    CROSS JOIN period_boundaries_cte pb\n    WHERE d.week_end_date <= pb.r4w_end_date\n      AND d.week_end_date >= pb.p4w_start_date\n),\nmomentum_base_cte AS (\n    SELECT\n        SUM(CASE WHEN d.week_end_date BETWEEN pb.r4w_start_date AND pb.r4w_end_date THEN d.relmora_total_mg ELSE 0 END) AS relmora_total_mg_r4w,\n        SUM(CASE WHEN d.week_end_date BETWEEN pb.p4w_start_date AND pb.p4w_end_date THEN d.relmora_total_mg ELSE 0 END) AS relmora_total_mg_p4w\n    FROM data_867 d\n    CROSS JOIN period_boundaries_cte pb\n    WHERE d.week_end_date <= pb.r4w_end_date\n      AND d.week_end_date >= pb.p4w_start_date\n),\nmomentum_metrics_incomplete_cte AS (\n    SELECT\n        pb.r4w_start_date,\n        pb.r4w_end_date,\n        pb.p4w_start_date,\n        pb.p4w_end_date,\n        bd.r4w_business_days,\n        bd.p4w_business_days,\n        pc.is_r4w_complete,\n        pc.is_p4w_complete,\n        mb.relmora_total_mg_r4w,\n        mb.relmora_total_mg_p4w,\n        CAST(mb.relmora_total_mg_r4w / NULLIF(bd.r4w_business_days, 0) AS INTEGER) AS relmora_daily_avg_mg_r4w,\n        CAST(mb.relmora_total_mg_p4w / NULLIF(bd.p4w_business_days, 0) AS INTEGER) AS relmora_daily_avg_mg_p4w,\n        CONCAT(\n            ROUND(\n                CASE\n                    WHEN CAST(mb.relmora_total_mg_p4w / NULLIF(bd.p4w_business_days, 0) AS INTEGER) = 0\n                         AND CAST(mb.relmora_total_mg_r4w / NULLIF(bd.r4w_business_days, 0) AS INTEGER) > 0 THEN 100\n                    WHEN CAST(mb.relmora_total_mg_p4w / NULLIF(bd.p4w_business_days, 0) AS INTEGER) = 0\n                         AND CAST(mb.relmora_total_mg_r4w / NULLIF(bd.r4w_business_days, 0) AS INTEGER) = 0 THEN 0\n                    ELSE (\n                        (\n                            CAST(mb.relmora_total_mg_r4w / NULLIF(bd.r4w_business_days, 0) AS INTEGER) -\n                            CAST(mb.relmora_total_mg_p4w / NULLIF(bd.p4w_business_days, 0) AS INTEGER)\n                        ) / NULLIF(CAST(mb.relmora_total_mg_p4w / NULLIF(bd.p4w_business_days, 0) AS INTEGER), 0)\n                    ) * 100\n                END\n            ),\n            '%'\n        ) AS relmora_daily_avg_growth_mg_r4w_vs_p4w\n    FROM period_boundaries_cte pb\n    CROSS JOIN business_day_counts_cte bd\n    CROSS JOIN period_completeness_cte pc\n    CROSS JOIN momentum_base_cte mb\n),\nmomentum_classification_cte AS (\n    SELECT\n        m.*,\n        CASE\n            WHEN ROUND(\n                CASE\n                    WHEN m.relmora_daily_avg_mg_p4w = 0 AND m.relmora_daily_avg_mg_r4w > 0 THEN 100\n                    WHEN m.relmora_daily_avg_mg_p4w = 0 AND m.relmora_daily_avg_mg_r4w = 0 THEN 0\n                    ELSE ((m.relmora_daily_avg_mg_r4w - m.relmora_daily_avg_mg_p4w) / NULLIF(m.relmora_daily_avg_mg_p4w, 0)) * 100\n                END\n            ) > 0 THEN 'Positive'\n            WHEN ROUND(\n                CASE\n                    WHEN m.relmora_daily_avg_mg_p4w = 0 AND m.relmora_daily_avg_mg_r4w > 0 THEN 100\n                    WHEN m.relmora_daily_avg_mg_p4w = 0 AND m.relmora_daily_avg_mg_r4w = 0 THEN 0\n                    ELSE ((m.relmora_daily_avg_mg_r4w - m.relmora_daily_avg_mg_p4w) / NULLIF(m.relmora_daily_avg_mg_p4w, 0)) * 100\n                END\n            ) < 0 THEN 'Negative'\n            ELSE 'Flat'\n        END AS short_term_sales_momentum_assessment_r4w_vs_p4w\n    FROM momentum_metrics_incomplete_cte m\n)\nSELECT\n    r4w_start_date,\n    r4w_end_date,\n    p4w_start_date,\n    p4w_end_date,\n    r4w_business_days,\n    p4w_business_days,\n    is_r4w_complete,\n    is_p4w_complete,\n    relmora_total_mg_r4w,\n    relmora_total_mg_p4w,\n    relmora_daily_avg_mg_r4w,\n    relmora_daily_avg_mg_p4w,\n    relmora_daily_avg_growth_mg_r4w_vs_p4w,\n    short_term_sales_momentum_assessment_r4w_vs_p4w\nFROM momentum_classification_cte\nLIMIT 1;\n\nResult Summary:\nThe analysis answers a straightforward business question: is Relmora’s national sales momentum improving or weakening in the most recent 4 weeks versus the prior 4 weeks. It looks at national Relmora sales volume over the latest eight-week span in the data, comparing Feb 7, 2026–Mar 6, 2026 against Jan 10, 2026–Feb 6, 2026. Both periods are complete, which means the comparison is reliable and can be read as a true like-for-like momentum check.\n\nThe result shows Relmora’s short-term momentum is negative. Total volume fell from 370,125 mg in the prior 4 weeks to 259,205 mg in the recent 4 weeks, and the daily average also declined from 19,480 mg to 17,280 mg, a drop of 11%. What stands out is that sales weakened even after adjusting for business days, indicating this is not just a calendar effect but a real slowdown in underlying sales pace.\n\nRelevant Questions:\n- Is short-term sales performance consistent across regions?\n- How are we doing in terms of adding new businesses?\n- What is the sales in the recent week?\n- In which regions is relmora gaining or losing share?", additional_kwargs={}, response_metadata={}, id='9dc9ff8b-906e-4e49-b700-34a5a4aadf95', tool_calls=[], invalid_tool_calls=[]), AIMessage(content='SQL query results', additional_kwargs={'type': 'sql_result', 'data': {'columns': ['R4W_START_DATE', 'R4W_END_DATE', 'P4W_START_DATE', 'P4W_END_DATE', 'R4W_BUSINESS_DAYS', 'P4W_BUSINESS_DAYS', 'IS_R4W_COMPLETE', 'IS_P4W_COMPLETE', 'RELMORA_TOTAL_MG_R4W', 'RELMORA_TOTAL_MG_P4W', 'RELMORA_DAILY_AVG_MG_R4W', 'RELMORA_DAILY_AVG_MG_P4W', 'RELMORA_DAILY_AVG_GROWTH_MG_R4W_VS_P4W', 'SHORT_TERM_SALES_MOMENTUM_ASSESSMENT_R4W_VS_P4W'], 'data': [{'R4W_START_DATE': datetime.date(2026, 2, 7), 'R4W_END_DATE': datetime.date(2026, 3, 6), 'P4W_START_DATE': datetime.date(2026, 1, 10), 'P4W_END_DATE': datetime.date(2026, 2, 6), 'R4W_BUSINESS_DAYS': 15, 'P4W_BUSINESS_DAYS': 19, 'IS_R4W_COMPLETE': 1, 'IS_P4W_COMPLETE': 1, 'RELMORA_TOTAL_MG_R4W': 259205, 'RELMORA_TOTAL_MG_P4W': 370125, 'RELMORA_DAILY_AVG_MG_R4W': 17280, 'RELMORA_DAILY_AVG_MG_P4W': 19480, 'RELMORA_DAILY_AVG_GROWTH_MG_R4W_VS_P4W': '-11%', 'SHORT_TERM_SALES_MOMENTUM_ASSESSMENT_R4W_VS_P4W': 'Negative'}]}}, response_metadata={}, id='c36aaffe-b9a6-4827-9d73-bfb5a0f1f0c8', tool_calls=[], invalid_tool_calls=[]), AIMessage(content='Visualization', additional_kwargs={'type': 'visualization', 'code': 'NO_VISUALIZATION'}, response_metadata={}, id='4ce4a58f-1e0e-47a5-bb3a-7627edbcd34a', tool_calls=[], invalid_tool_calls=[])]

client = OpenAI()
def convert_to_dataframe(sql_data):
    if not sql_data:
        return None

    return pd.DataFrame(sql_data["data"], columns=sql_data["columns"])

def parse_conversation(messages):
    """
    Converts cleaned messages into structured blocks.

    Now includes:
    - question
    - summary
    - data
    - viz_code ✅
    """

    blocks = []
    current_block = None

    for msg in messages:

        # 🟢 Human Message → start new block
        if isinstance(msg, HumanMessage):
            if current_block:
                blocks.append(current_block)

            current_block = {
                "question": msg.content,
                "summary": None,
                "data": None,
                "viz_code": None   # ✅ NEW
            }

        # 🔵 AI Message → extract info
        elif isinstance(msg, AIMessage) and current_block:

            content = msg.content or ""
            kwargs = getattr(msg, "additional_kwargs", {}) or {}
            msg_type = kwargs.get("type")

            # ✅ 1. SQL Result
            if msg_type == "sql_result":
                current_block["data"] = kwargs.get("data")
                continue

            # ✅ 2. Visualization (THIS IS WHAT YOU WANT)
            if msg_type == "visualization":
                code = kwargs.get("code")

                # Handle NO_VISUALIZATION case
                if code and code != "NO_VISUALIZATION":
                    current_block["viz_code"] = code
                else:
                    current_block["viz_code"] = None

                continue

            # ✅ 3. Extract Summary
            if "Result Summary:" in content:
                summary = content.split("Result Summary:")[-1].strip()
                current_block["summary"] = summary
                continue

            # ⚠️ Fallback
            if current_block["summary"] is None and len(content) > 100:
                current_block["summary"] = content.strip()

    # Append last block
    if current_block:
        blocks.append(current_block)

    return blocks

def generate_slide_content(block):
    """
    Uses LLM to convert block into slide-ready content.

    Input:
        block = {
            "question": "...",
            "summary": "...",
            "data": {...}
        }

    Output:
        {
            "title": "...",
            "bullets": ["...", "..."],
            "kpis": ["...", "..."]
        }
    """

    question = block.get("question", "")
    summary = block.get("summary", "")
    data_sample = str(block.get("data"))[:1000]  # truncate for token safety

    prompt = f"""

You are a commercial operations consultant preparing a single PPT slide 
for senior leadership at a life sciences firm.
 
Your goal is to highlight growth, demand signals, risks, opportunities, 
and strategic implications — NOT just describe data.
 
You MUST adapt the output based on BOTH the DATA DOMAIN and the 
ANALYTICAL PATTERN. These are independent and must be classified separately.
 
------------------------------------------------------------
INPUT:
Question: {question}
Summary: {summary}
Data Sample: {data_sample}
------------------------------------------------------------
 
### STEP 1: CLASSIFY DATA DOMAIN (mandatory first step)
 
Read the question, summary, and data sample. Identify which domain the 
underlying metrics belong to. Use the lexicon below to decide.
 
DOMAIN A — SALES_PERFORMANCE
  Signals: sales trends, sales performance, breadth, depth, adding new businesses, account adoption, own-brand performance
  Vocabulary allowed: sales, growth, demand, volume, units, new accounts, new business
 
DOMAIN B — COMPETITOR_DYNAMICS  
  Signals:  market share, competitor, gaining/lossing share, competitor share, 
  Vocabulary allowed: share, share-shift, competitive position,  capture rate
 
DOMAIN C — FIELD_EXECUTION
  Signals: calls, call activity, touch, touch activity, reach, frequency, rep visits, coverage of HCPs/accounts by reps
  Vocabulary allowed: calls, reach, frequency, coverage, call plan, touches, effort percentage
 
DOMAIN D — CROSS_DOMAIN
    Signals: capture rate, sales growth comparison  across our product and competitor products or among competitors, 
  Use ONLY when the question explicitly correlates two domains, 
  e.g., "did increased calls drive sales lift?" or "How does our product sales grow in recent 3 months and how does that compare with competitor sales?"
 
CRITICAL RULE: 
If the question is FIELD_EXECUTION, you MUST NOT use sales/revenue/share 
vocabulary in title, bullets, KPIs, or insight.
If the question is SALES_PERFORMANCE, you MUST NOT use call/reach/
frequency vocabulary unless the data sample contains it.
If the question is COMPETITOR_DYNAMICS, you MUST NOT report own-brand 
sales as the headline metric.
For cross-domain questions, the context will inherently span multiple domains such as FIELD_EXECUTION, SALES_PERFORMANCE, and/or COMPETITOR_DYNAMICS. In such scenarios, it is acceptable to use vocabulary from multiple relevant domains together, while ensuring the terminology used in titles, bullets, KPIs, and insights remains contextually appropriate and aligned with the dominant signal in the data.
If the data sample contradicts the question's apparent domain, trust 
the data sample and flag the mismatch in the insight.
 
### STEP 2: CLASSIFY ANALYTICAL PATTERN
 
[Your existing 10 categories — TREND, SHORT_TERM, REGIONAL, DISTRIBUTION, 
ACCOUNT_HEALTH, NEW_BUSINESS, ADOPTION, TIER, MARKET_SHARE, TOP_N]
 
Note: MARKET_SHARE pattern is only valid with COMPETITOR_DYNAMICS domain.
 
### STEP 3: SELECT KPIs
 
Select KPIs only from the "Findings" section of the provided summary and strictly anchor 
all KPI generation to that content. Generate a maximum of 3 KPIs only.

── KPI ORDERING & STRUCTURE ───────────────────────────────────────────────────

- The FIRST KPI must always represent National Sales and Growth combined, and 
  must follow this exact format:
  "National [Product] [Period] Sales: [prior sales] → [current sales] ([growth%])"
  
    ✅ VALID:   "National Relmora R4W Sales: 2.4M → 2.1M (-30%)"
    ✅ VALID:   "National Account Metric: 135 → 129 (-4%)"
    ❌ INVALID: Showing growth % alone without sales values in the first KPI.
    ❌ INVALID: Using a geography or tier metric as the first KPI even if it 
                shows a larger movement than the national metric.

- The first KPI MUST combine both:
    • Absolute sales values (prior → current)
    • Growth percentage in brackets alongside
  These two MUST always appear together in the first KPI — never separately.

- If both National Sales and National Growth are unavailable, apply the 
  following fallback priority order for the first KPI:
    1. National Sales alone (if only sales is available):
         ✅ "National Relmora R4W Sales: 2.4M → 2.1M"
    2. National Growth alone (if only growth % is available):
         ✅ "National Relmora R4W Growth: -30%"
    3. If neither National Sales nor National Growth is present in the 
       findings — only then may the first KPI slot be filled with the 
       most significant geography-level metric available.
    ❌ NEVER leave the first KPI slot empty if any national metric exists.

- The SECOND and THIRD KPIs must represent the most significant 
  geography-level or tier-level trends/metrics from the findings, in 
  descending order of significance.

── KPI VALUE RULES ────────────────────────────────────────────────────────────

- Each KPI must contain a clearly defined label and the exact corresponding 
  value directly supported by the findings.
- Preserve all numeric comparisons exactly as stated, including current vs 
  prior values, growth changes, and absolute sales/volume figures.
- Do not infer, introduce, summarize, or calculate KPIs beyond what is 
  explicitly stated in the findings.
- If National Sales or growth data is not present in the findings, the first 
  KPI slot must remain empty — do NOT substitute a geography metric in its place.

── REFERENCE EXAMPLES ─────────────────────────────────────────────────────────

    National Relmora R4W Sales:   2.4M → 2.1M (-30%)
    National Account Growth:      18% → 16%
    National R3M Sales:           2.4M units
    Regions Above Benchmark:      West, Central, Great Lakes
    Regions Below Benchmark:      South East

### STEP 4: INSIGHT GENERATION LOGIC
INSIGHT_GENERATION → Form insights strictly from the "Opportunity / Implication" section of the provided summary, with primary emphasis on highlighting the opportunity areas and business implications. Ensure all insights are directly anchored to the provided content without introducing unsupported assumptions or conclusions.
 

### STEP 5: BULLET GENERATION LOGIC
BULLET_GENERATION → Form bullets strictly from the "Key Takeaways" section of the provided summary. Ensure every bullet is directly derived from the provided key takeaways without introducing any new interpretation, assumption, metric, or conclusion. Preserve the original meaning, comparisons, entity standings, and notable gaps or contrasts exactly as described in the summary, Focus on the metrics and display it.

### STEP 6: KPI Defination
Given a KPI label and value, write a one-line definition (max 15 words) 
explaining what the metric measures. Be specific, avoid filler words.

------------------------------------------------------------

### OUTPUT FORMAT EXAMPLE (STRICT JSON ONLY):

{{
    "title": "How does new account addition look like across regions?",
    "bullets": [
        "National new-account additions declined 4% versus the prior period.",
        "West, Central, and Great Lakes outperformed the national benchmark.",
        "West emerged as the strongest region, leading both in volume and relative performance."
    ],
    "kpis": [
        {{"label": "National Account Growth", "value": "-4%","defination":"Percentage change in new accounts added versus the prior period, nationally."}},
        {{"label": "Regions Above Benchmark", "value": "West, Central, Great Lakes", "defination":"Regions whose new-account growth rate exceeded the national average this period."}},
        {{"label": "Regions Below Benchmark", "value": "North East, South East, Mid Atlantic","defination":"Regions whose new-account growth rate fell short of the national average this period."}}
    ],
    "insight": "National new-account additions declined 4%, indicating softer acquisition momentum overall. West, Central, and Great Lakes outperformed the national benchmark and can serve as reference markets for successful acquisition strategies. The largest opportunity lies in improving performance across North East, South East, and Mid Atlantic, which are contributing most to the overall regional drag."
}}

------------------------------------------------------------

### STRICT GUIDELINES

TITLE:
- Max 20 words
- Grammatically correct version of: {question}
- Preserve original meaning without adding new interpretation
- Do not convert into an insight or statement beyond correcting grammar

BULLETS:
- Exactly 3 bullets


KPIs:
- Max 3 KPIs

INSIGHT:
- 2–3 lines max
- Must connect data → implication → action

------------------------------------------------------------

### IMPORTANT BEHAVIOR RULES

- DO NOT repeat summary text
- DO NOT describe charts/data literally
- ALWAYS interpret patterns
- ALWAYS highlight:
  → Growth OR Risk OR Opportunity
- Prefer sharp, executive language

------------------------------------------------------------

### FINAL SELF-CHECK

✔ Title ≤ 20 words  
✔ Exactly 3 bullets  
✔ Insight is actionable  
✔ KPIs match question type  
✔ No generic statements  
✔ Output is leadership-ready  
Return ONLY valid JSON. No extra text.
"""

    response = client.chat.completions.create(
        model="gpt-5.4",
        messages=[{"role": "user", "content": prompt}],
        temperature=0
    )

    content = response.choices[0].message.content

    try:
        return json.loads(content)
    except:
        print("⚠️ JSON parsing failed, raw output:", content)
        return {
            "title": question,
            "bullets": [summary[:100]],
            "kpis": []
        }
    
def build_slide_object(block):
    slide_llm = generate_slide_content(block)

    return {
        "title": slide_llm["title"],
        "bullets": slide_llm["bullets"],
        "kpis": slide_llm["kpis"],
        "insight":slide_llm["insight"],
        "data": block.get("data"),
        "viz_code": block.get("viz_code")
    }

def _normalize_plotly_figure(figure: dict) -> dict | None:
    try:
        return json.loads(json.dumps(figure, cls=PlotlyJSONEncoder))
    except Exception:
        return None


def generate_chart(block):
    if not block.get("viz_code"):
        return None

    df = convert_to_dataframe(block["data"])

    local_vars = {}

    try:
        exec(block["viz_code"], {"df": df}, local_vars)

        fig = local_vars.get("fig")
        if fig is None:
            return None

        file_path = f"chart_{uuid.uuid4()}.png"
        fig.write_image(file_path)

        return file_path

    except Exception as e:
        print("Chart error:", e)
        return None
    

def enrich_with_chart(slide):
    chart_path = generate_chart({
        "data": slide["data"],
        "viz_code": slide["viz_code"]
    })

    slide["chart_path"] = chart_path
    return slide  


def build_slide_data(
    messages,
    chart_path_overrides: list[str] | None = None,
    cancel_check: Callable[[], None] | None = None,
):
    blocks = parse_conversation(messages)
    print(f"[PPT] slides: building (blocks={len(blocks)})")
    slides = []
    for index, block in enumerate(blocks):
        if cancel_check:
            cancel_check()
        slide = build_slide_object(block)
        if isinstance(chart_path_overrides, list):
            override = chart_path_overrides[index] if index < len(chart_path_overrides) else None
            slide["chart_path"] = override
        else:
            slide = enrich_with_chart(slide, cancel_check=cancel_check)
        slides.append(slide)
    print(f"[PPT] slides: built (slides={len(slides)})")
    return slides






import os
import zipfile
import xml.etree.ElementTree as ET
from datetime import date
from typing import Any
 
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
 
 
# ──────────────────────────────────────────────────────────────────────────────
# DEFAULT PALETTE  (white-background safe)
# ──────────────────────────────────────────────────────────────────────────────
 
_DEFAULT = dict(
    background     = RGBColor(0xFF, 0xFF, 0xFF),
    title_text     = RGBColor(0x0D, 0x1B, 0x3E),
    body_text      = RGBColor(0x0D, 0x1B, 0x3E),
    muted_text     = RGBColor(0x55, 0x6B, 0x88),
    accent1        = RGBColor(0x00, 0xB4, 0xD8),
    accent2        = RGBColor(0x00, 0xE5, 0xFF),
    accent3        = RGBColor(0x39, 0xD3, 0x53),
    panel_bg       = RGBColor(0xF0, 0xF4, 0xF8),
    kpi_bg         = RGBColor(0xF7, 0xF9, 0xFC),
    kpi_label_text = RGBColor(0x1E, 0x6E, 0xA8),
    kpi_border     = RGBColor(0xCC, 0xD6, 0xE0),
    subtitle_label = RGBColor(0x1E, 0x6E, 0xA8),
    bullet_dot     = RGBColor(0x00, 0xB4, 0xD8),
    divider        = RGBColor(0x00, 0xB4, 0xD8),
    insight_label  = RGBColor(0x00, 0xB4, 0xD8),
    insight_text   = RGBColor(0x55, 0x6B, 0x88),
    font_title     = "Calibri",
    font_body      = "Calibri",
)
 
# DrawingML namespace
_DML = "http://schemas.openxmlformats.org/drawingml/2006/main"
 
 
# ──────────────────────────────────────────────────────────────────────────────
# COLOR / FONT HELPERS
# ──────────────────────────────────────────────────────────────────────────────
 
def _hex_to_rgb(hex_str: str) -> RGBColor | None:
    try:
        h = hex_str.lstrip("#")
        if len(h) != 6:
            return None
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return None
 
 
def _luminance(rgb: RGBColor) -> float:
    return 0.2126 * rgb[0] / 255 + 0.7152 * rgb[1] / 255 + 0.0722 * rgb[2] / 255
 
 
def _lighten(rgb: RGBColor, factor: float = 0.85) -> RGBColor:
    return RGBColor(
        int(rgb[0] + (255 - rgb[0]) * factor),
        int(rgb[1] + (255 - rgb[1]) * factor),
        int(rgb[2] + (255 - rgb[2]) * factor),
    )
 
 
def _mix(c1: RGBColor, c2: RGBColor, t: float = 0.5) -> RGBColor:
    return RGBColor(
        int(c1[0] + (c2[0] - c1[0]) * t),
        int(c1[1] + (c2[1] - c1[1]) * t),
        int(c1[2] + (c2[2] - c1[2]) * t),
    )
 
 
# ──────────────────────────────────────────────────────────────────────────────
# THEME EXTRACTOR  ← THE FIXED FUNCTION
# ──────────────────────────────────────────────────────────────────────────────
 
def extract_theme_from_pptx(path: str) -> dict:
    """
    Extract brand theme (colors + fonts) from a .pptx file and return a dict
    ready for use by create_ppt().
 
    Strategy
    --------
    1. Open the .pptx as a zip and parse ppt/theme/theme*.xml with ElementTree.
       This bypasses python-pptx's SlideMaster object entirely, which is the
       source of the 'has no attribute theme_color_map' error.
    2. Scan slide XML files for explicit typeface attributes to find the dominant
       run-level font (fallback when theme XML uses '+mn-lt' placeholders).
    3. Map raw theme slots → semantic roles used by the slide builder.
 
    Always returns a complete dict (falls back to _DEFAULT for missing values).
    """
    theme = dict(_DEFAULT)
 
    try:
        with zipfile.ZipFile(path, "r") as z:
            all_files = z.namelist()
 
            # ── 1. Parse theme XML files ───────────────────────────────────
            theme_files = sorted(
                f for f in all_files if f.startswith("ppt/theme/theme") and f.endswith(".xml")
            )
 
            raw_colors: dict[str, RGBColor] = {}
            font_title: str | None = None
            font_body:  str | None = None
 
            for tf_path in theme_files:
                xml_bytes = z.read(tf_path)
                root = ET.fromstring(xml_bytes)
 
                # ── Color scheme ──────────────────────────────────────────
                clr_scheme = root.find(f".//{{{_DML}}}clrScheme")
                if clr_scheme is not None and not raw_colors:
                    for child in clr_scheme:
                        slot = child.tag.split("}")[-1]          # e.g. "dk1", "accent1"
                        # srgbClr (hex color)
                        srgb = child.find(f"{{{_DML}}}srgbClr")
                        if srgb is not None:
                            c = _hex_to_rgb(srgb.get("val", ""))
                            if c:
                                raw_colors[slot] = c
                        # sysClr (system color with lastClr fallback)
                        sys_clr = child.find(f"{{{_DML}}}sysClr")
                        if sys_clr is not None:
                            c = _hex_to_rgb(sys_clr.get("lastClr", ""))
                            if c:
                                raw_colors[slot] = c
 
                # ── Font scheme ───────────────────────────────────────────
                font_scheme = root.find(f".//{{{_DML}}}fontScheme")
                if font_scheme is not None and font_title is None:
                    major = font_scheme.find(f"{{{_DML}}}majorFont/{{{_DML}}}latin")
                    minor = font_scheme.find(f"{{{_DML}}}minorFont/{{{_DML}}}latin")
                    if major is not None:
                        tf = major.get("typeface", "")
                        # Skip theme-reference placeholders like "+mj-lt"
                        if tf and not tf.startswith("+"):
                            font_title = tf
                    if minor is not None:
                        tf = minor.get("typeface", "")
                        if tf and not tf.startswith("+"):
                            font_body = tf
 
            # ── 2. Scan slide XML for dominant explicit typeface ───────────
            slide_files = [f for f in all_files if f.startswith("ppt/slides/slide") and f.endswith(".xml")]
            font_counts: dict[str, int] = {}
 
            for sf in slide_files[:10]:     # sample first 10 slides for speed
                xml_bytes = z.read(sf)
                root = ET.fromstring(xml_bytes)
                for rPr in root.iter(f"{{{_DML}}}rPr"):
                    latin = rPr.find(f"{{{_DML}}}latin")
                    if latin is not None:
                        tf = latin.get("typeface", "")
                        if tf and not tf.startswith("+"):
                            font_counts[tf] = font_counts.get(tf, 0) + 1
 
            # Also scan slide master
            master_files = [f for f in all_files if "slideMasters/slideMaster" in f and f.endswith(".xml")]
            for mf in master_files:
                xml_bytes = z.read(mf)
                root = ET.fromstring(xml_bytes)
                for rPr in root.iter(f"{{{_DML}}}rPr"):
                    latin = rPr.find(f"{{{_DML}}}latin")
                    if latin is not None:
                        tf = latin.get("typeface", "")
                        if tf and not tf.startswith("+"):
                            font_counts[tf] = font_counts.get(tf, 0) + 1
 
            # Use dominant run font as body font if theme XML had placeholders
            if font_counts:
                dom_font = max(font_counts, key=font_counts.get)
                if font_body is None:
                    font_body = dom_font
                if font_title is None:
                    font_title = dom_font
 
        # ── 3. Map raw slots → semantic roles ──────────────────────────────
        accent_keys = ["accent1", "accent2", "accent3", "accent4", "accent5", "accent6"]
        accents = [raw_colors[k] for k in accent_keys if k in raw_colors]
 
        dk1 = raw_colors.get("dk1")
        dk2 = raw_colors.get("dk2")
 
        primary_accent   = accents[0] if len(accents) > 0 else _DEFAULT["accent1"]
        secondary_accent = accents[1] if len(accents) > 1 else _DEFAULT["accent2"]
        tertiary_accent  = accents[2] if len(accents) > 2 else _DEFAULT["accent3"]
 
        # Text: use dk1 if it's dark enough for a white background
        text_color = (dk1 if dk1 and _luminance(dk1) < 0.5
                      else _DEFAULT["title_text"])
        body_color = (dk2 if dk2 and _luminance(dk2) < 0.5
                      else text_color)
 
        panel_bg   = _lighten(primary_accent, 0.92)
        kpi_bg     = _lighten(primary_accent, 0.95)
        kpi_border = _lighten(primary_accent, 0.75)
        muted      = _mix(text_color, RGBColor(0xFF, 0xFF, 0xFF), 0.55)
        kpi_label  = _mix(primary_accent, text_color, 0.35)
 
        theme.update(
            background     = RGBColor(0xFF, 0xFF, 0xFF),   # always white
            title_text     = text_color,
            body_text      = body_color,
            muted_text     = muted,
            accent1        = primary_accent,
            accent2        = secondary_accent,
            accent3        = tertiary_accent,
            panel_bg       = panel_bg,
            kpi_bg         = kpi_bg,
            kpi_label_text = kpi_label,
            kpi_border     = kpi_border,
            subtitle_label = primary_accent,
            bullet_dot     = primary_accent,
            divider        = primary_accent,
            insight_label  = primary_accent,
            insight_text   = muted,
            font_title     = font_title or _DEFAULT["font_title"],
            font_body      = font_body  or _DEFAULT["font_body"],
        )
 
        print(f"[Theme] Extracted from '{path}':")
        print(f"  dk1={raw_colors.get('dk1')}  accent1={raw_colors.get('accent1')}")
        print(f"  font_title='{theme['font_title']}'  font_body='{theme['font_body']}'")
 
    except Exception as exc:
        print(f"[WARNING] Could not extract theme from '{path}': {exc}")
        print("[WARNING] Falling back to default theme.")
 
    return theme
 
 
# ──────────────────────────────────────────────────────────────────────────────
# LOW-LEVEL DRAWING HELPERS
# ──────────────────────────────────────────────────────────────────────────────
 
def _solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
 
def _no_line(shape):
    shape.line.fill.background()
 
def _rect(slide, l, t, w, h, color: RGBColor, border: bool = False):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    _solid(s, color)
    if not border:
        _no_line(s)
    return s
 
def _tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
 
def _para(tb, text: str, size: float, bold=False,
          color: RGBColor = None, italic=False,
          align=PP_ALIGN.LEFT, wrap=True,
          font: str = "Calibri", spacing=None,
          theme: dict = None):
    if color is None and theme:
        color = theme["body_text"]
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    if spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))
    return tb
 
def _label(slide, text: str, l, t, w, h,
           size=7.5, color: RGBColor = None, bold=False,
           align=PP_ALIGN.LEFT, italic=False,
           spacing=None, theme: dict = None):
    if color is None and theme:
        color = theme["muted_text"]
    tb = _tb(slide, l, t, w, h)
    return _para(tb, text, size, bold=bold, color=color, italic=italic,
                 align=align, wrap=False, spacing=spacing)
 
 
# ──────────────────────────────────────────────────────────────────────────────
# KPI CARD
# ──────────────────────────────────────────────────────────────────────────────
 
def _parse_kpi(kpi: Any) -> tuple[str, str]:
    if isinstance(kpi, dict):
        return str(kpi.get("label", "")).strip(), str(kpi.get("value", "")).strip()
    parts = str(kpi).split(":", 1)
    return parts[0].strip(), (parts[1].strip() if len(parts) > 1 else "")
 
 
def _kpi_card(slide, label: str, value: str,
              l, t, w, h, accent_color: RGBColor, theme: dict):
    card = _rect(slide, l, t, w, h, theme["kpi_bg"])
    card.line.color.rgb = theme["kpi_border"]
    card.line.width = Pt(0.75)
 
    _rect(slide, l, t, 0.06, h, accent_color)   # left accent bar
 
    inner_l = l + 0.14
    inner_w = w - 0.18
 
    # Label
    tb_lbl = _tb(slide, inner_l, t + 0.10, inner_w, 0.25)
    tf_lbl = tb_lbl.text_frame
    tf_lbl.word_wrap = True
    p_lbl = tf_lbl.paragraphs[0]
    p_lbl.alignment = PP_ALIGN.LEFT
    run_lbl = p_lbl.add_run()
    run_lbl.text = label
    run_lbl.font.size = Pt(7.5)
    run_lbl.font.bold = True
    run_lbl.font.name = theme["font_body"]
    run_lbl.font.color.rgb = theme["kpi_label_text"]
    rPr = run_lbl._r.get_or_add_rPr()
    rPr.set("spc", str(int(2.2 * 100)))
 
    # Value
    tb = _tb(slide, inner_l, t + 0.33, inner_w, h - 0.38)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = value
    run.font.size = Pt(18 if len(value) > 12 else 22)
    run.font.bold = True
    run.font.name = theme["font_title"]
    run.font.color.rgb = theme["title_text"]
 
 
# ──────────────────────────────────────────────────────────────────────────────
# BULLETS PANEL
# ──────────────────────────────────────────────────────────────────────────────
 
def _bullets_panel(slide, bullets: list[str], l, t, w, h, theme: dict):
    panel = _rect(slide, l, t, w, h, theme["panel_bg"])
    panel.line.color.rgb = theme["kpi_border"]
    panel.line.width = Pt(0.75)
 
    if not bullets:
        tb = _tb(slide, l + 0.3, t + h / 2 - 0.22, w - 0.6, 0.44)
        _para(tb, "No insights available.", 13,
              italic=True, color=theme["muted_text"], align=PP_ALIGN.CENTER)
        return
 
    n = len(bullets)
    PAD_L, PAD_T, PAD_B, DOT_W, ROW_GAP = 0.22, 0.18, 0.14, 0.22, 0.10
    avail_h = h - PAD_T - PAD_B
    row_h   = (avail_h - ROW_GAP * (n - 1)) / n
 
    for i, bullet_text in enumerate(bullets):
        row_y = t + PAD_T + i * (row_h + ROW_GAP)
 
        dot_tb = _tb(slide, l + PAD_L, row_y + row_h * 0.18, DOT_W, row_h)
        dot_tf = dot_tb.text_frame
        dot_tf.word_wrap = False
        dot_p  = dot_tf.paragraphs[0]
        dot_run = dot_p.add_run()
        dot_run.text = "●"
        dot_run.font.size = Pt(8)
        dot_run.font.color.rgb = theme["bullet_dot"]
        dot_run.font.name = theme["font_body"]
 
        text_l = l + PAD_L + DOT_W
        text_w = w - PAD_L - DOT_W - 0.15
        tb = _tb(slide, text_l, row_y, text_w, row_h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet_text
        run.font.size = Pt(13)
        run.font.bold = False
        run.font.name = theme["font_body"]
        run.font.color.rgb = theme["body_text"]
 
 
# ──────────────────────────────────────────────────────────────────────────────
# MAIN SLIDE BUILDER
# ──────────────────────────────────────────────────────────────────────────────
 
def create_ppt(slide_data: dict,
               prs: Presentation | None = None,
               theme: dict | None = None,
               logo_path: str | None = None) -> Presentation:
    if theme is None:
        theme = dict(_DEFAULT)
 
    if prs is None:
        prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
 
    slide = prs.slides.add_slide(prs.slide_layouts[6])
 
    has_chart   = bool(slide_data.get("chart_path") and os.path.exists(slide_data["chart_path"]))
    has_kpis    = bool(slide_data.get("kpis"))
    has_insight = bool(slide_data.get("insight"))
    bullets     = slide_data.get("bullets") or []
 
    W, H = 10.0, 5.625
 
    # Background (always white)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["background"]
 
    MARGIN_L, MARGIN_R, MARGIN_T = 0.38, 0.30, 0.18
    CONTENT_L  = MARGIN_L
    SUBTITLE_H = 0.24
    TITLE_H    = 0.48
    DIVIDER_Y  = MARGIN_T + SUBTITLE_H + TITLE_H + 0.04
    DIVIDER_H  = 0.03
    FOOTER_H   = 0.22
    FOOTER_Y   = H - FOOTER_H - 0.06
 
    if has_insight:
        INSIGHT_H   = 0.52
        INSIGHT_Y   = FOOTER_Y - INSIGHT_H - 0.04
        CONTENT_BOT = INSIGHT_Y - 0.10
    else:
        CONTENT_BOT = FOOTER_Y - 0.06
 
    CONTENT_T = DIVIDER_Y + DIVIDER_H + 0.08
    CONTENT_H = CONTENT_BOT - CONTENT_T
    CONTENT_W = W - MARGIN_L - MARGIN_R
 
    KPI_GAP = 0.18
    if has_kpis:
        KPI_W        = 2.85
        LEFT_PANEL_W = CONTENT_W - KPI_W - KPI_GAP
        LEFT_PANEL_X = CONTENT_L
        KPI_X        = CONTENT_L + LEFT_PANEL_W + KPI_GAP
    else:
        KPI_W = LEFT_PANEL_X = KPI_X = 0
        LEFT_PANEL_W = CONTENT_W
        LEFT_PANEL_X = CONTENT_L
 
    # Subtitle
    subtitle = slide_data.get("subtitle", "BUSINESS QUESTION")
    _label(slide, subtitle, CONTENT_L, MARGIN_T, 5, SUBTITLE_H,
           size=8, color=theme["subtitle_label"], bold=True, spacing=2.5)
 
    # Title
    tb = _tb(slide, CONTENT_L, MARGIN_T + SUBTITLE_H, W - MARGIN_L - MARGIN_R, TITLE_H)
    _para(tb, slide_data.get("title", "Untitled"),
          24, bold=True, color=theme["title_text"], font=theme["font_title"])
 
    # Divider
    _rect(slide, CONTENT_L, DIVIDER_Y, CONTENT_W, DIVIDER_H, theme["divider"])
 
    # Chart or bullets panel
    if has_chart:
        panel = _rect(slide, LEFT_PANEL_X, CONTENT_T, LEFT_PANEL_W, CONTENT_H, theme["panel_bg"])
        panel.line.color.rgb = theme["kpi_border"]
        panel.line.width = Pt(0.75)
        PAD = 0.06
        slide.shapes.add_picture(
            slide_data["chart_path"],
            Inches(LEFT_PANEL_X + PAD), Inches(CONTENT_T + PAD),
            Inches(LEFT_PANEL_W - PAD * 2), Inches(CONTENT_H - PAD * 2)
        )
    else:
        _bullets_panel(slide, bullets, LEFT_PANEL_X, CONTENT_T, LEFT_PANEL_W, CONTENT_H, theme=theme)
 
    # KPI cards
    # if has_kpis:
    #     kpis          = slide_data["kpis"][:3]
    #     n             = len(kpis)
    #     GAP           = 0.14
    #     card_h        = (CONTENT_H - GAP * (n - 1)) / n
    #     accent_colors = [theme["accent1"], theme["accent2"], theme["accent3"]]
    #     for i, kpi in enumerate(kpis):
    #         cy           = CONTENT_T + i * (card_h + GAP)
    #         label, value = _parse_kpi(kpi)
    #         _kpi_card(slide, label, value, KPI_X, cy, KPI_W, card_h,
    #                   accent_colors[i % len(accent_colors)], theme=theme)

    # KPI cards
    if has_kpis:
        kpis          = slide_data["kpis"][:3]
        n             = len(kpis)
        GAP           = 0.14
        card_h        = (CONTENT_H - GAP * (n - 1)) / n
        accent_colors = [theme["accent1"], theme["accent2"], theme["accent3"]]
        for i, kpi in enumerate(kpis):
            cy           = CONTENT_T + i * (card_h + GAP)
            label, value = _parse_kpi(kpi)
            definition   = kpi.get("defination", "") if isinstance(kpi, dict) else ""
            _kpi_card(slide, label, value, KPI_X, cy, KPI_W, card_h,
                      accent_colors[i % len(accent_colors)], theme=theme)

            # Definition text box — pinned to bottom of each card
            if definition:
                DEF_H   = 0.22
                DEF_PAD = 0.06
                tb = _tb(slide,
                         KPI_X + DEF_PAD,
                         cy + card_h - DEF_H - 0.04,
                         KPI_W - DEF_PAD * 2,
                         DEF_H)
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = definition
                run.font.size = Pt(6.5)
                run.font.color.rgb = theme.get("muted_text", RGBColor(0x99, 0x99, 0x99))
                run.font.name = theme["font_body"]
                run.font.italic = True
 
   
    # Key Insight
    if has_insight:
        KI_LABEL_W = 1.05
        KI_LABEL_H = 0.22   # 👈 FIX: explicit height for label

        # Label (now constrained height)
        _label(slide, "KEY INSIGHT",
            CONTENT_L, INSIGHT_Y,
            KI_LABEL_W, KI_LABEL_H,
            size=8.5, color=theme["insight_label"], bold=True, spacing=1.5)

        # --- LOGO BELOW LABEL (no overlap now) ---
        LOGO_W = 0.8
        LOGO_H = 0.4
        LOGO_MARGIN = 0.2

        if logo_path and os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path,
                Inches(CONTENT_L),
                Inches(INSIGHT_Y + KI_LABEL_H + LOGO_MARGIN),
                width=Inches(LOGO_W),
                height=Inches(LOGO_H)
            )

        # --- INSIGHT TEXT (UNCHANGED POSITION) ---
        tb = _tb(slide,
                CONTENT_L + KI_LABEL_W + 0.06,
                INSIGHT_Y,
                CONTENT_W - KI_LABEL_W - 0.06,
                INSIGHT_H)

        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        run.text = slide_data["insight"]
        run.font.size = Pt(9)
        run.font.color.rgb = theme["insight_text"]
        run.font.name = theme["font_body"]
 
    # Footer
    footnote     = slide_data.get("footnote", "")
    today        = date.today().strftime("%d %b %Y")
    FOOTER_COLOR = theme["muted_text"]
 
    if footnote:
        _label(slide, footnote, CONTENT_L, FOOTER_Y, CONTENT_W * 0.65, FOOTER_H,
               size=7, color=FOOTER_COLOR, align=PP_ALIGN.CENTER)
    # else:
    #     _label(slide, f"Generated {today}", CONTENT_L, FOOTER_Y, 3, FOOTER_H,
    #            size=7, color=FOOTER_COLOR)
 
    # _label(slide, "CONFIDENTIAL  ·  FOR INTERNAL USE ONLY",
    #        W - MARGIN_R - 3.2, FOOTER_Y, 3.2, FOOTER_H,
    #        size=7, color=FOOTER_COLOR, align=PP_ALIGN.RIGHT, spacing=0.5)
 
    return prs
 
 
# ──────────────────────────────────────────────────────────────────────────────
# CONVENIENCE BUILDER
# ──────────────────────────────────────────────────────────────────────────────
 
def build_ppt(
    messages,
    output_path: str = "final_presentation.pptx",
    uploaded_pptx_path: str | None = "Geron.pptx",
    logo_path: str | None = "Geron_Logo.png",
    chart_path_overrides: list[str] | None = None,
    cancel_check: Callable[[], None] | None = None,
) -> str:
    resolved_template = (
        uploaded_pptx_path
        if uploaded_pptx_path and os.path.exists(uploaded_pptx_path)
        else None
    )
    resolved_logo = logo_path if logo_path and os.path.exists(logo_path) else None
    theme = (extract_theme_from_pptx(resolved_template)
             if resolved_template else dict(_DEFAULT))

    print(
        "[PPT] build: start "
        f"output={output_path} template={resolved_template or ''} logo={resolved_logo or ''}"
    )
 
    if cancel_check:
        cancel_check()
    slides = build_slide_data(
        messages,
        chart_path_overrides=chart_path_overrides,
        cancel_check=cancel_check,
    )
 
    prs = None
    for slide in slides:
        if cancel_check:
            cancel_check()
        prs = create_ppt(slide, prs, theme=theme, logo_path=resolved_logo)
 
    if cancel_check:
        cancel_check()
    prs.save(output_path)
    print(f"[PPT] build: saved output={output_path} slides={len(slides)}")

    for slide in slides:
        chart_path = slide.get("chart_path")
        if chart_path and os.path.exists(chart_path):
            try:
                os.remove(chart_path)
            except Exception:
                print(f"[PPT] cleanup: failed to remove chart {chart_path}")
    return output_path


#build_ppt(messages)